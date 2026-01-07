"""
Study Pilot AI - Document Ingestion Pipeline
Parse PDFs, PowerPoints, and text files for RAG
"""

import re
import json
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from dataclasses import dataclass, asdict

# PDF parsing
try:
    from PyPDF2 import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

# PowerPoint parsing
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


@dataclass
class DocumentChunk:
    """A chunk of document content for embedding."""
    id: str
    content: str
    source_file: str
    source_type: str  # 'pdf', 'pptx', 'txt', 'syllabus'
    page_or_slide: int
    course_id: Optional[int]
    topic_id: Optional[int]
    metadata: Dict


class DocumentIngester:
    """
    Ingests various document types and creates chunks for RAG.
    Supports: PDF, PowerPoint, plain text, JSON syllabi
    """
    
    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 100):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.chunk_counter = 0
    
    def ingest_file(self, file_path: Path, course_id: int = None) -> List[DocumentChunk]:
        """
        Ingest a single file based on its extension.
        
        Args:
            file_path: Path to the file
            course_id: Optional course ID to associate
            
        Returns:
            List of DocumentChunk objects
        """
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        
        if suffix == '.pdf':
            return self.ingest_pdf(file_path, course_id)
        elif suffix in ['.pptx', '.ppt']:
            return self.ingest_pptx(file_path, course_id)
        elif suffix == '.txt':
            return self.ingest_text(file_path, course_id)
        elif suffix == '.json':
            return self.ingest_json_syllabus(file_path, course_id)
        else:
            print(f"Unsupported file type: {suffix}")
            return []
    
    def ingest_pdf(self, file_path: Path, course_id: int = None) -> List[DocumentChunk]:
        """Extract and chunk content from PDF."""
        if not HAS_PYPDF:
            print("PyPDF2 not installed. Cannot process PDFs.")
            return []
        
        chunks = []
        
        try:
            reader = PdfReader(str(file_path))
            
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                text = self._clean_text(text)
                
                if not text.strip():
                    continue
                
                # Create chunks from this page
                page_chunks = self._create_chunks(
                    text=text,
                    source_file=file_path.name,
                    source_type='pdf',
                    page_or_slide=page_num,
                    course_id=course_id
                )
                chunks.extend(page_chunks)
        
        except Exception as e:
            print(f"Error processing PDF {file_path}: {e}")
        
        return chunks
    
    def ingest_pptx(self, file_path: Path, course_id: int = None) -> List[DocumentChunk]:
        """Extract and chunk content from PowerPoint."""
        if not HAS_PPTX:
            print("python-pptx not installed. Cannot process PowerPoints.")
            return []
        
        chunks = []
        
        try:
            prs = Presentation(str(file_path))
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                text = "\n".join(slide_text)
                text = self._clean_text(text)
                
                if not text.strip():
                    continue
                
                # Create chunks from this slide
                slide_chunks = self._create_chunks(
                    text=text,
                    source_file=file_path.name,
                    source_type='pptx',
                    page_or_slide=slide_num,
                    course_id=course_id
                )
                chunks.extend(slide_chunks)
        
        except Exception as e:
            print(f"Error processing PPTX {file_path}: {e}")
        
        return chunks
    
    def ingest_text(self, file_path: Path, course_id: int = None) -> List[DocumentChunk]:
        """Ingest plain text file."""
        chunks = []
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            text = self._clean_text(text)
            
            if text.strip():
                chunks = self._create_chunks(
                    text=text,
                    source_file=file_path.name,
                    source_type='txt',
                    page_or_slide=1,
                    course_id=course_id
                )
        
        except Exception as e:
            print(f"Error processing text file {file_path}: {e}")
        
        return chunks
    
    def ingest_json_syllabus(self, file_path: Path, 
                             course_id: int = None) -> List[DocumentChunk]:
        """Ingest structured syllabus from JSON."""
        chunks = []
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Extract course description
            if 'description' in data:
                chunks.extend(self._create_chunks(
                    text=data['description'],
                    source_file=file_path.name,
                    source_type='syllabus',
                    page_or_slide=0,
                    course_id=course_id
                ))
            
            # Extract topics
            for i, topic in enumerate(data.get('topics', []), 1):
                topic_text = f"Topic: {topic.get('name', '')}\n"
                topic_text += f"Description: {topic.get('description', '')}\n"
                
                if 'subtopics' in topic:
                    topic_text += f"Subtopics: {', '.join(topic['subtopics'])}\n"
                
                if 'key_concepts' in topic:
                    topic_text += f"Key Concepts: {', '.join(topic['key_concepts'])}\n"
                
                chunks.extend(self._create_chunks(
                    text=topic_text,
                    source_file=file_path.name,
                    source_type='syllabus',
                    page_or_slide=i,
                    course_id=course_id,
                    topic_id=topic.get('id')
                ))
            
            # Extract from questions if present (for context)
            for q in data.get('questions', [])[:10]:  # Limit
                q_text = f"Q: {q.get('question_text', '')}\n"
                q_text += f"Answer: {q.get('correct_answer', '')}\n"
                q_text += f"Explanation: {q.get('explanation', '')}"
                
                chunks.extend(self._create_chunks(
                    text=q_text,
                    source_file=file_path.name,
                    source_type='syllabus',
                    page_or_slide=q.get('topic_id', 0),
                    course_id=course_id,
                    topic_id=q.get('topic_id')
                ))
        
        except Exception as e:
            print(f"Error processing JSON syllabus {file_path}: {e}")
        
        return chunks
    
    def ingest_directory(self, dir_path: Path, 
                         course_id: int = None) -> List[DocumentChunk]:
        """Ingest all supported files in a directory."""
        dir_path = Path(dir_path)
        all_chunks = []
        
        supported_extensions = ['.pdf', '.pptx', '.ppt', '.txt', '.json']
        
        for file_path in dir_path.iterdir():
            if file_path.suffix.lower() in supported_extensions:
                chunks = self.ingest_file(file_path, course_id)
                all_chunks.extend(chunks)
                print(f"Ingested {file_path.name}: {len(chunks)} chunks")
        
        return all_chunks
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special characters but keep basic punctuation
        text = re.sub(r'[^\w\s.,;:!?\-\(\)\[\]\'\"]+', '', text)
        return text.strip()
    
    def _create_chunks(self, text: str, source_file: str, source_type: str,
                       page_or_slide: int, course_id: int = None,
                       topic_id: int = None) -> List[DocumentChunk]:
        """Split text into overlapping chunks."""
        chunks = []
        
        if len(text) <= self.chunk_size:
            # Single chunk
            self.chunk_counter += 1
            chunks.append(DocumentChunk(
                id=f"chunk_{self.chunk_counter}",
                content=text,
                source_file=source_file,
                source_type=source_type,
                page_or_slide=page_or_slide,
                course_id=course_id,
                topic_id=topic_id,
                metadata={
                    'char_count': len(text),
                    'word_count': len(text.split())
                }
            ))
        else:
            # Split into overlapping chunks
            start = 0
            while start < len(text):
                end = start + self.chunk_size
                
                # Try to end at a sentence boundary
                if end < len(text):
                    # Look for sentence end within last 100 chars
                    search_zone = text[end-100:end]
                    for sep in ['. ', '! ', '? ', '\n']:
                        last_sep = search_zone.rfind(sep)
                        if last_sep != -1:
                            end = end - 100 + last_sep + len(sep)
                            break
                
                chunk_text = text[start:end].strip()
                
                if chunk_text:
                    self.chunk_counter += 1
                    chunks.append(DocumentChunk(
                        id=f"chunk_{self.chunk_counter}",
                        content=chunk_text,
                        source_file=source_file,
                        source_type=source_type,
                        page_or_slide=page_or_slide,
                        course_id=course_id,
                        topic_id=topic_id,
                        metadata={
                            'char_count': len(chunk_text),
                            'word_count': len(chunk_text.split()),
                            'chunk_index': len(chunks)
                        }
                    ))
                
                start = end - self.chunk_overlap
        
        return chunks
    
    def chunks_to_dict(self, chunks: List[DocumentChunk]) -> List[dict]:
        """Convert chunks to dictionaries for storage."""
        return [asdict(chunk) for chunk in chunks]
